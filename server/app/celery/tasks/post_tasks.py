"""
Module for processing uploaded files asynchronously and extracting text embeddings.

This module defines a Celery task `process_uploaded_file` that
asynchronously processes uploaded files, extracts text content,
and generates text embeddings. It also provides utility functions
for decoding base64 encoded strings to ObjectId and extracting text
from PDF files.

Tasks:
    process_uploaded_file: Asynchronously processes an uploaded file,
    extracts text content, and generates text embeddings.

Utilities:
    extract_text_from_pdf: Extracts text content from a PDF file.
    extract_text_embedding: Generates text embeddings from text content.

"""

import mimetypes
import io
import os
import math
from uuid import UUID

from flask import current_app
import fitz
from app.celery.celery import celery_instance
from celery.signals import task_success, task_failure
import google.generativeai as genai
from app.models.embedding import Embedding
from mongoengine import connect
from dotenv import load_dotenv
from pptx import Presentation
from docx import Document


def extract_text_from_pdf(file_data: bytes) -> str:
    """
    Extract text content from a PDF file.

    Args:
        file_data (bytes): The binary data of the PDF file.

    Returns:
        str: The extracted text content from the PDF.

    Raises:
        Exception: If an error occurs during PDF processing.

    """
    text = ""
    try:
        # Open the PDF file
        pdf_document = fitz.open(stream=io.BytesIO(file_data), filetype="pdf")

        # Iterate through each page in the PDF
        for page_number in range(len(pdf_document)):
            page = pdf_document.load_page(page_number)

            # Extract text from the page
            text += page.get_text()

        return text

    except Exception as error:
        print(f"Error: {error}")
        raise


def extract_text_from_ppt(file_data: bytes) -> str:
    """
    Extract text from a PowerPoint (PPT) file given its byte data.

    This function reads a PowerPoint (PPT) file from the provided byte data and
    extracts all available text content from each slide. It iterates through each
    slide in the presentation, examining each shape on the slide. If a shape contains
    text, the text content is extracted and concatenated into a single string.

    Args:
        file_data (bytes): Byte data representing the PowerPoint (PPT) file.

    Returns:
        str: The extracted text content from the PowerPoint (PPT) file.

    Raises:
        Exception: If an error occurs during the extraction process.

    Note:
        - The function uses the python-pptx library to handle PowerPoint file parsing
          and text extraction.
        - The provided byte data should represent a valid PowerPoint file (.pptx).
        - Shapes that contain text, such as text boxes or shapes with textual content,
          are considered when extracting text.
        - Any non-textual shapes, such as images or graphical elements without text
          content, are not considered for text extraction.
        - If an error occurs during the extraction process, an exception is raised
          with details about the error.
    """
    try:
        ppt_stream = io.BytesIO(file_data)
        presentation = Presentation(ppt_stream)

        extracted_text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"

        return extracted_text

    except Exception as error:
        print(f"Error: {error}")
        raise


def extract_text_from_docx(file_data: bytes) -> str:
    """
    Extract text from a Word document (docx file) given its byte data.

    This function reads a Word document (docx file) from the provided byte data and
    extracts all available text content. It iterates through each paragraph in the
    document and extracts the text content from each paragraph. The extracted text
    includes content from paragraphs, headings, titles, and lists.

    Args:
        file_data (bytes): Byte data representing the Word document (docx file).

    Returns:
        str: The extracted text content from the Word document (docx file).

    Raises:
        Exception: If an error occurs during the extraction process.

    Note:
        - The function uses the python-docx library to handle Word document parsing
          and text extraction.
        - The provided byte data should represent a valid Word document file (.docx).
        - The extracted text includes content from paragraphs, headings, titles, and
          lists present in the document.
        - Each paragraph in the document is separated by a newline character ('\n') in
          the extracted text.
        - If an error occurs during the extraction process, an exception is raised with
          details about the error.
    """
    try:
        docx_stream = io.BytesIO(file_data)
        doc = Document(docx_stream)

        extracted_text = ""

        for paragraph in doc.paragraphs:
            extracted_text += paragraph.text + "\n"

        return extracted_text

    except Exception as error:
        print(f"Error: {error}")
        raise


def extract_text_embedding(chunk: str) -> list:
    """
    Generate text embeddings for a text chunk using a pre-trained model.

    Args:
        chunk (str): The text chunk for which embeddings are to be generated.

    Returns:
        list: A list of embedding vectors representing the text chunk.

    Raises:
        Exception: If an error occurs during the embedding generation process.

    """
    try:
        result = genai.embed_content(
            model="models/embedding-001",
            content=chunk,
            task_type="semantic_similarity",
        )
        return result["embedding"]
    except Exception as error:
        print(f"Error: {error}")
        raise


@celery_instance.task(soft_time_limit=60, time_limit=120)
def process_uploaded_file(
    file_data: bytes,
    filename: str,
    hub_id: str,
    post_id: UUID,
    attachment_id: str,
) -> None:
    """
    Asynchronously process an uploaded file, extract text content,
    generate text embeddings, and save them to MongoDB.

    Args:
        file_data (bytes): The binary data of the uploaded file.
        filename (str): The name of the uploaded file.
        hub_id (str): The ID of the hub to which the file belongs.
        post_id (UUID): The UUID of the post to which the file belongs.

    Returns:
        None

    Raises:
        Exception: If an error occurs during file processing or embedding generation.

    """
    try:
        load_dotenv()
        connect(
            db=os.getenv("MONGO_DB"),
            host=os.getenv("MONGO_URI"),
            username=os.getenv("MONGO_USERNAME"),
            password=os.getenv("MONGO_PASSWORD"),
            alias="default",
        )
        print("Connected to MongoDB successfully!")

        file_type = mimetypes.guess_type(filename)[0]
        extracted_text = ""

        if file_type == "application/pdf":
            extracted_text = extract_text_from_pdf(file_data)
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            extracted_text = extract_text_from_ppt(file_data)
        elif (
            file_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            extracted_text = extract_text_from_docx(file_data)
        else:
            raise ValueError(
                "Invalid file format. Only PDF, PPT, and DOCX formats are supported."
            )

        embedding_docs = []

        num_chunks = len(extracted_text)
        counter = 0

        for i in range(0, num_chunks, 1000):
            chunk = extracted_text[i : i + 1000]
            embedding = extract_text_embedding(chunk)
            counter += 1
            embedding_doc = Embedding(
                hub_id=hub_id,
                post_id=post_id,
                attachment_id=attachment_id,
                batch_no=counter,
                text_content=chunk,
                embeddings=embedding,
            )
            embedding_docs.append(embedding_doc)

        Embedding.objects.insert(embedding_docs, load_bulk=False)
        attachment_number_of_embeddings_key = (
            f"attachment_id_{attachment_id}_number_of_embeddings"
        )

        redis_client = current_app.redis_client

        redis_client.set(
            attachment_number_of_embeddings_key, math.ceil(num_chunks / 1000)
        )

    except Exception as error:
        print(f"error: {error}")


@task_success.connect(sender=process_uploaded_file)
def task_success_handler(sender=None, result=None, **kwargs):
    """
    Event handler function triggered when a Celery task succeeds.

    Args:
        sender: The sender of the signal (the Celery task).
        result: The result returned by the Celery task.
        **kwargs: Additional keyword arguments.

    Returns:
        None

    """
    task_id = sender.request.id
    print("SUCCESS")
    redis_client = current_app.redis_client
    redis_client.publish(f"{task_id}", "SUCCESS")


@task_failure.connect(sender=process_uploaded_file)
def task_failure_handler(sender=None, exception=None, traceback=None, **kwargs):
    """
    Event handler function triggered when a Celery task fails.

    Args:
        sender: The sender of the signal (the Celery task).
        exception: The exception raised by the Celery task.
        traceback: The traceback associated with the exception.
        **kwargs: Additional keyword arguments.

    Returns:
        None

    """
    task_id = sender.request.id
    print("FAILURE")
    redis_client = current_app.redis_client
    redis_client.publish(f"{task_id}", "FAILURE")
